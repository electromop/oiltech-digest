"""Разбор загруженного файла в список якорей.

Чистые функции: ни базы, ни сети, ни обращений к модели. На вход — байты файла,
на выход — ParsedDocument из model.py. Всё, что не разобралось, поднимается
как DocumentError с русским текстом, который можно показать пользователю.

Формат определяется ПО СИГНАТУРЕ, а не по имени файла: имя приходит от браузера
и ничего не гарантирует (переименованный .txt в .pdf — обычное дело).

Библиотеки выбраны замером и зафиксированы: pypdf (35 МБ пик против 511 МБ
у pdfplumber; pymupdf отпал по AGPL), python-pptx, python-docx.
"""

from __future__ import annotations

import io
import zipfile

from docx import Document as _DocxDocument
from docx.table import Table as _DocxTable
from docx.text.paragraph import Paragraph as _DocxParagraph
from pptx import Presentation as _Presentation
from pypdf import PdfReader

from .model import (
    ANCHOR_UNITS,
    Anchor,
    ParsedDocument,
    ScannedDocumentError,
    UnreadableDocumentError,
    UnsupportedFormatError,
)

__all__ = ["detect_kind", "parse_document", "looks_scanned"]


# --- сигнатуры -------------------------------------------------------------

_PDF_SIGNATURE = b"%PDF-"
_ZIP_SIGNATURE = b"PK\x03\x04"
# Старые бинарные форматы Office и файлы, защищённые паролем на открытие:
# и то и другое — контейнер OLE2, снаружи неотличимы без сторонней библиотеки.
_OLE2_SIGNATURE = b"\xd0\xcf\x11\xe0\xa1\xb1\x1a\xe1"

# Внутренние имена внутри OOXML-архива. Различают pptx и docx: снаружи оба — zip.
_PPTX_MARKER = "ppt/presentation.xml"
_DOCX_MARKER = "word/document.xml"
_XLSX_MARKER = "xl/workbook.xml"

# Заголовок %PDF- по стандарту стоит в начале файла, но встречается мусорный
# префикс (почтовые шлюзы, конвертеры). pypdf такие файлы читает, поэтому ищем
# сигнатуру в первом килобайте, а не строго с нулевого байта.
_PDF_HEADER_WINDOW = 1024


# --- пороги «это скан» -----------------------------------------------------
#
# Порог считается по ОБЪЁМУ добытого текста, а не по доле пустых якорей.
# Замер, который это доказал: «Росатом_позиция_+_заключение_экспертизы.pdf» —
# 107 страниц, текстовый слой есть только на 40 (37%), но это 104 431 знак
# полезного текста. Любой критерий вида «пустых страниц больше половины —
# отклоняем» выбросил бы этот документ целиком. Он обязан разбираться.
#
# Настоящий скан выглядит иначе: текста нет НИГДЕ. Извлекатель либо возвращает
# пустоту, либо крохи вроде колонцифры на каждой странице.

# Якорь считается содержательным, если в нём есть хотя бы короткая строка.
# 20 знаков — примерно «Отчёт за 2026 год»: меньше этого встречается у
# колонтитулов и номеров страниц, то есть у артефактов скана.
SCAN_MEANINGFUL_ANCHOR_CHARS = 20

# Второй критерий — плотность: если якорей много, а знаков на якорь почти нет,
# это скан с редкими вкраплениями текста (водяной знак, штамп, колонцифра).
# Нужен, потому что первого критерия мало: одна страница с текстом среди сотни
# пустых иначе прошла бы как полноценный документ.
SCAN_MIN_DENSITY = 5.0
SCAN_DENSITY_MIN_ANCHORS = 10


def detect_kind(data: bytes) -> str:
    """Определить формат по сигнатуре. Вернуть "pdf" | "pptx" | "docx"."""
    if not isinstance(data, (bytes, bytearray, memoryview)):
        raise TypeError("detect_kind ожидает байты файла")
    data = bytes(data)

    if not data:
        raise UnsupportedFormatError("Файл пустой — загрузите PDF, PPTX или DOCX.")

    if _PDF_SIGNATURE in data[:_PDF_HEADER_WINDOW]:
        return "pdf"

    if data.startswith(_ZIP_SIGNATURE):
        return _detect_ooxml_kind(data)

    if data.startswith(_OLE2_SIGNATURE):
        raise UnsupportedFormatError(
            "Похоже на старый формат Office (.doc, .ppt, .xls) или на файл, "
            "защищённый паролем на открытие. Пересохраните его как PDF, PPTX "
            "или DOCX и загрузите снова."
        )

    raise UnsupportedFormatError(
        "Не удалось распознать формат файла. Поддерживаются PDF, PPTX и DOCX."
    )


def _detect_ooxml_kind(data: bytes) -> str:
    """Разобрать zip-контейнер: pptx и docx снаружи одинаковы (оба PK\\x03\\x04)."""
    try:
        with zipfile.ZipFile(io.BytesIO(data)) as archive:
            names = set(archive.namelist())
    except zipfile.BadZipFile as exc:
        # Сигнатура zip есть, а содержимое не читается — это битый файл,
        # а не «неизвестный формат»: пользователю надо перезалить, а не конвертировать.
        raise UnreadableDocumentError(
            "Файл повреждён: архив документа не читается. Попробуйте пересохранить "
            "и загрузить снова."
        ) from exc

    if _PPTX_MARKER in names:
        return "pptx"
    if _DOCX_MARKER in names:
        return "docx"
    if _XLSX_MARKER in names:
        raise UnsupportedFormatError(
            "Это таблица Excel. Поддерживаются PDF, PPTX и DOCX."
        )
    raise UnsupportedFormatError(
        "Внутри архива нет ни презентации, ни документа Word. "
        "Поддерживаются PDF, PPTX и DOCX."
    )


def parse_document(data: bytes, kind: str | None = None) -> ParsedDocument:
    """Разобрать файл в якоря. kind=None — определить формат самостоятельно."""
    if not isinstance(data, (bytes, bytearray, memoryview)):
        raise TypeError("parse_document ожидает байты файла")
    data = bytes(data)

    if kind is None:
        kind = detect_kind(data)
    elif kind not in ANCHOR_UNITS:
        raise UnsupportedFormatError(
            f"Формат «{kind}» не поддерживается. Поддерживаются PDF, PPTX и DOCX."
        )

    if kind == "pdf":
        anchors = _parse_pdf(data)
    elif kind == "pptx":
        anchors = _parse_pptx(data)
    else:
        anchors = _parse_docx(data)

    document = ParsedDocument(kind=kind, anchors=anchors)
    if looks_scanned(anchors, kind):
        unit = ANCHOR_UNITS[kind]
        raise ScannedDocumentError(
            f"В документе нет текстового слоя — похоже, это скан. "
            f"Проверено единиц («{unit}»): {len(anchors)}, извлечено знаков: "
            f"{document.total_chars}. Распознавание изображений пока не "
            f"поддерживается — загрузите файл с текстом."
        )
    return document


# --- PDF -------------------------------------------------------------------


def _parse_pdf(data: bytes) -> list[Anchor]:
    """Якорь = страница, нумерация 1-based. Пустые страницы СОХРАНЯЮТСЯ."""
    try:
        reader = PdfReader(io.BytesIO(data), strict=False)
        if reader.is_encrypted:
            # Часть PDF зашифрована пустым паролем (запрет печати/копирования) —
            # такие читаются. Настоящий пароль на открытие вернёт ложь.
            try:
                opened = reader.decrypt("")
            except Exception:  # noqa: BLE001 — pypdf кидает разное на битой криптографии
                opened = 0
            if not opened:
                raise UnreadableDocumentError(
                    "PDF защищён паролем. Снимите пароль и загрузите снова."
                )
        page_count = len(reader.pages)
    except UnreadableDocumentError:
        raise
    except Exception as exc:  # noqa: BLE001 — pypdf кидает десяток разных классов
        raise UnreadableDocumentError(
            "Не удалось прочитать PDF: файл повреждён. Попробуйте пересохранить "
            "и загрузить снова."
        ) from exc

    anchors: list[Anchor] = []
    for number in range(1, page_count + 1):
        try:
            text = reader.pages[number - 1].extract_text() or ""
        except Exception:  # noqa: BLE001
            # Одна битая страница не должна ронять весь документ: на 107-страничном
            # файле из-за неё пропали бы 106 нормальных страниц.
            text = ""
        anchors.append(Anchor(number=number, text=text.strip()))
    return anchors


# --- PPTX ------------------------------------------------------------------


def _parse_pptx(data: bytes) -> list[Anchor]:
    """Якорь = слайд. Берём фигуры с текстом, таблицы и заметки докладчика."""
    try:
        presentation = _Presentation(io.BytesIO(data))
        slides = list(presentation.slides)
    except Exception as exc:  # noqa: BLE001
        raise UnreadableDocumentError(
            "Не удалось прочитать презентацию: файл повреждён или защищён паролем."
        ) from exc

    anchors: list[Anchor] = []
    for number, slide in enumerate(slides, start=1):
        parts: list[str] = []
        try:
            _collect_shapes(slide.shapes, parts)
            notes = _slide_notes(slide)
            if notes:
                parts.append(notes)
        except Exception:  # noqa: BLE001 — битый слайд не роняет презентацию
            pass
        text = "\n".join(part.strip() for part in parts if part and part.strip())
        anchors.append(Anchor(number=number, text=text.strip()))
    return anchors


def _collect_shapes(shapes, parts: list[str]) -> None:
    """Обойти фигуры слайда рекурсивно: в реальных колодах текст лежит в группах."""
    for shape in shapes:
        try:
            if hasattr(shape, "shapes"):  # группа фигур
                _collect_shapes(shape.shapes, parts)
                continue
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    parts.append(" | ".join(cell.text.strip() for cell in row.cells))
                continue
            if getattr(shape, "has_text_frame", False):
                parts.append(shape.text_frame.text)
        except Exception:  # noqa: BLE001 — одна фигура не роняет слайд
            continue


def _slide_notes(slide) -> str:
    if not getattr(slide, "has_notes_slide", False):
        return ""
    frame = slide.notes_slide.notes_text_frame
    return frame.text if frame is not None else ""


# --- DOCX ------------------------------------------------------------------


def _parse_docx(data: bytes) -> list[Anchor]:
    """Якорь = порядковый номер НЕПУСТОГО блока.

    Страниц в docx нет — разбиение на страницы делает Word при вёрстке, в файле
    его не хранится. Поэтому нумеруем блоки: абзац или таблица целиком.
    Идём по телу документа в исходном порядке, иначе таблицы уехали бы в конец
    и ссылка «блок 12» перестала бы соответствовать глазу читателя.
    """
    try:
        document = _DocxDocument(io.BytesIO(data))
        body = document.element.body
    except Exception as exc:  # noqa: BLE001
        raise UnreadableDocumentError(
            "Не удалось прочитать документ Word: файл повреждён или защищён паролем."
        ) from exc

    anchors: list[Anchor] = []
    for child in body.iterchildren():
        tag = child.tag.rsplit("}", 1)[-1]
        try:
            if tag == "p":
                text = _DocxParagraph(child, document).text
            elif tag == "tbl":
                text = _table_text(_DocxTable(child, document))
            else:
                continue
        except Exception:  # noqa: BLE001 — битый блок пропускаем, документ читаем
            continue
        text = text.strip()
        if not text:
            continue  # пустые абзацы-разделители якорями не считаем
        anchors.append(Anchor(number=len(anchors) + 1, text=text))
    return anchors


def _table_text(table: _DocxTable) -> str:
    rows = []
    for row in table.rows:
        rows.append(" | ".join(cell.text.strip() for cell in row.cells))
    return "\n".join(rows)


# --- критерий скана --------------------------------------------------------


def looks_scanned(anchors: list[Anchor], kind: str) -> bool:
    """Скан ли это. См. обоснование порогов в шапке модуля.

    Документ с РЕДКИМ текстом сканом не считается: критерий смотрит на объём
    добытого текста, а не на долю пустых якорей.
    """
    if not anchors:
        return True

    meaningful = sum(
        1 for anchor in anchors if len(anchor.text.strip()) >= SCAN_MEANINGFUL_ANCHOR_CHARS
    )
    if meaningful == 0:
        # Нигде в документе не нашлось даже одной короткой строки.
        return True

    # Плотность меряем только там, где якорь — физическая страница или слайд.
    # У docx пустые блоки в якоря не попадают вовсе, там эта метрика бессмысленна.
    if kind in ("pdf", "pptx") and len(anchors) >= SCAN_DENSITY_MIN_ANCHORS:
        total_chars = sum(len(anchor.text) for anchor in anchors)
        if total_chars / len(anchors) < SCAN_MIN_DENSITY:
            return True

    return False
