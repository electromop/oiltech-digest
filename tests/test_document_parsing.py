"""Тесты разбора загруженных документов.

Основная часть работает на синтетике, собранной прямо здесь: набор должен
проходить на любой машине. Калибровочные тесты на реальных файлах владельца
помечены skipif — без файлов на диске они пропускаются, а не падают.
"""

from __future__ import annotations

import io
import os
import zipfile

import pytest
from docx import Document as DocxDocument
from pptx import Presentation
from pptx.util import Inches
from pypdf import PdfWriter

from oiltech_digest.documents.model import (
    Anchor,
    ScannedDocumentError,
    UnreadableDocumentError,
    UnsupportedFormatError,
)
from oiltech_digest.documents.parsing import (
    SCAN_MEANINGFUL_ANCHOR_CHARS,
    SCAN_MIN_DENSITY,
    detect_kind,
    looks_scanned,
    parse_document,
)


# --- сборка синтетических файлов ------------------------------------------


def _content_stream(text: str) -> bytes:
    if not text:
        return b"BT ET"  # страница есть, текста на ней нет — имитация скана
    escaped = text.replace("\\", r"\\").replace("(", r"\(").replace(")", r"\)")
    return f"BT /F1 12 Tf 72 720 Td ({escaped}) Tj ET".encode("latin-1")


def make_pdf(pages: list[str]) -> bytes:
    """Минимальный валидный PDF с текстовым слоем.

    Текст только латиницей: базовый Helvetica без встроенной кодировки отдаёт
    кириллицу мусором, а проверяем мы парсер, а не шрифты.
    """
    objects: dict[int, bytes] = {
        3: b"<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>",
    }
    next_id = 4
    page_ids: list[int] = []
    for text in pages:
        stream = _content_stream(text)
        content_id = next_id
        next_id += 1
        objects[content_id] = b"<< /Length %d >>\nstream\n%s\nendstream" % (
            len(stream),
            stream,
        )
        page_id = next_id
        next_id += 1
        objects[page_id] = (
            b"<< /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792] "
            b"/Resources << /Font << /F1 3 0 R >> >> /Contents %d 0 R >>" % content_id
        )
        page_ids.append(page_id)

    kids = b" ".join(b"%d 0 R" % pid for pid in page_ids)
    objects[2] = b"<< /Type /Pages /Kids [%s] /Count %d >>" % (kids, len(page_ids))
    objects[1] = b"<< /Type /Catalog /Pages 2 0 R >>"

    out = bytearray(b"%PDF-1.4\n")
    offsets: dict[int, int] = {}
    for oid in sorted(objects):
        offsets[oid] = len(out)
        out += b"%d 0 obj\n" % oid + objects[oid] + b"\nendobj\n"

    xref_pos = len(out)
    size = max(objects) + 1
    out += b"xref\n0 %d\n" % size
    out += b"0000000000 65535 f \n"
    for oid in range(1, size):
        out += b"%010d 00000 n \n" % offsets[oid]
    out += b"trailer\n<< /Size %d /Root 1 0 R >>\nstartxref\n%d\n%%%%EOF\n" % (
        size,
        xref_pos,
    )
    return bytes(out)


def make_pptx(slides: list[dict]) -> bytes:
    """Презентация из описаний слайдов: text, table (список строк), notes."""
    presentation = Presentation()
    blank = presentation.slide_layouts[6]
    for spec in slides:
        slide = presentation.slides.add_slide(blank)
        if spec.get("text"):
            box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(1))
            box.text_frame.text = spec["text"]
        rows = spec.get("table")
        if rows:
            shape = slide.shapes.add_table(
                len(rows), len(rows[0]), Inches(1), Inches(3), Inches(6), Inches(2)
            )
            for r, row in enumerate(rows):
                for c, value in enumerate(row):
                    shape.table.cell(r, c).text = value
        if spec.get("notes"):
            slide.notes_slide.notes_text_frame.text = spec["notes"]
    buffer = io.BytesIO()
    presentation.save(buffer)
    return buffer.getvalue()


def make_docx(paragraphs: list[str], table: list[list[str]] | None = None) -> bytes:
    document = DocxDocument()
    for text in paragraphs:
        document.add_paragraph(text)
    if table:
        added = document.add_table(rows=len(table), cols=len(table[0]))
        for r, row in enumerate(table):
            for c, value in enumerate(row):
                added.cell(r, c).text = value
    buffer = io.BytesIO()
    document.save(buffer)
    return buffer.getvalue()


def make_plain_zip() -> bytes:
    buffer = io.BytesIO()
    with zipfile.ZipFile(buffer, "w") as archive:
        archive.writestr("readme.txt", "не офисный документ")
    return buffer.getvalue()


# --- detect_kind -----------------------------------------------------------


def test_detect_kind_pdf():
    assert detect_kind(make_pdf(["Hello, this is a page"])) == "pdf"


def test_detect_kind_separates_pptx_and_docx():
    # Оба формата — zip с одинаковой сигнатурой PK\x03\x04, различает только
    # содержимое архива. Тест ловит регрессию «всё zip считаем docx».
    pptx_bytes = make_pptx([{"text": "Slide"}])
    docx_bytes = make_docx(["Абзац"])
    assert pptx_bytes[:4] == docx_bytes[:4] == b"PK\x03\x04"
    assert detect_kind(pptx_bytes) == "pptx"
    assert detect_kind(docx_bytes) == "docx"


@pytest.mark.parametrize(
    "data",
    [
        pytest.param(b"", id="пусто"),
        pytest.param(b"\x00\x01\x02\x03\xff\xfe", id="мусор"),
        pytest.param(b"\x89PNG\r\n\x1a\n" + b"\x00" * 32, id="png"),
        pytest.param(b"PK\x05\x06" + b"\x00" * 18, id="пустой-архив"),
        pytest.param(b"%PDF is a nice format, but this is plain text", id="почти-pdf"),
    ],
)
def test_detect_kind_rejects_unknown_bytes(data):
    # Требование: мусор даёт понятную ошибку, а не падение парсера.
    with pytest.raises(UnsupportedFormatError):
        detect_kind(data)


def test_detect_kind_ignores_filename_and_looks_at_bytes():
    # Имя приходит от браузера и ничего не гарантирует: содержимое — текст,
    # хоть «файл и называется отчёт.pdf».
    with pytest.raises(UnsupportedFormatError):
        detect_kind("Это отчёт.pdf, честное слово".encode("utf-8"))


def test_detect_kind_rejects_legacy_office():
    ole2 = b"\xd0\xcf\x11\xe0\xa1\xb1\x1a\xe1" + b"\x00" * 64
    with pytest.raises(UnsupportedFormatError) as exc:
        detect_kind(ole2)
    assert "паролем" in str(exc.value) or ".doc" in str(exc.value)


def test_detect_kind_rejects_zip_without_office_parts():
    with pytest.raises(UnsupportedFormatError):
        detect_kind(make_plain_zip())


def test_detect_kind_reports_broken_archive_as_unreadable():
    # Сигнатура zip есть, содержимое обрезано — это битый файл, а не чужой формат.
    broken = make_docx(["Абзац"])[:200]
    with pytest.raises(UnreadableDocumentError):
        detect_kind(broken)


# --- parse_document: PDF ---------------------------------------------------


def test_parse_pdf_anchor_per_page():
    document = parse_document(
        make_pdf(
            [
                "First page of the oilfield services report",
                "Second page with the drilling numbers",
                "Third page with the conclusion of the report",
            ]
        )
    )
    assert document.kind == "pdf"
    assert document.anchor_unit == "страница"
    assert [a.number for a in document.anchors] == [1, 2, 3]
    assert "Second page" in document.anchors[1].text
    assert document.total_chars > 0


def test_parse_pdf_keeps_empty_pages_and_numbering():
    """Пустые страницы остаются якорями — иначе съедет нумерация.

    Это уменьшенная копия файла «Росатом»: текст есть не везде, но ссылка
    «страница 4» должна указывать на четвёртую страницу исходника.
    """
    document = parse_document(
        make_pdf(
            [
                "Alpha page with enough text to count",
                "",
                "",
                "Delta page with enough text to count",
            ]
        )
    )
    assert len(document.anchors) == 4
    assert [a.number for a in document.anchors] == [1, 2, 3, 4]
    assert document.anchors[1].text == ""
    assert document.anchors[2].text == ""
    assert "Delta" in document.anchors[3].text
    empty = [a.number for a in document.anchors if not a.text]
    assert empty == [2, 3]  # вызывающий может сказать, какие якоря пустые


def test_parse_pdf_broken_file_is_unreadable():
    with pytest.raises(UnreadableDocumentError):
        parse_document(b"%PDF-1.4\nthis is not really a pdf body\n%%EOF")


def _encrypt_pdf(data: bytes, user_password: str, owner_password: str = "owner") -> bytes:
    writer = PdfWriter(clone_from=io.BytesIO(data))
    writer.encrypt(user_password=user_password, owner_password=owner_password)
    buffer = io.BytesIO()
    writer.write(buffer)
    return buffer.getvalue()


def test_parse_pdf_with_open_password_is_unreadable():
    data = _encrypt_pdf(make_pdf(["Secret page with enough text to count"]), "s3cret")
    with pytest.raises(UnreadableDocumentError) as exc:
        parse_document(data)
    assert "паролем" in str(exc.value)


def test_parse_pdf_with_owner_password_only_is_readable():
    # Запрет печати/копирования шифрует файл пустым паролем на открытие —
    # такие документы читаются, отказывать в них нельзя.
    data = _encrypt_pdf(make_pdf(["Restricted but readable page of text"]), "")
    document = parse_document(data)
    assert len(document.anchors) == 1
    assert "Restricted" in document.anchors[0].text


def test_parse_pdf_without_text_layer_is_scanned():
    with pytest.raises(ScannedDocumentError):
        parse_document(make_pdf(["", "", ""]))


# --- parse_document: PPTX --------------------------------------------------


def test_parse_pptx_collects_shapes_tables_and_notes():
    data = make_pptx(
        [
            {"text": "Заголовок первого слайда про нефтесервис"},
            {
                "text": "Второй слайд",
                "table": [["Показатель", "Значение"], ["Выручка", "1 200 млн"]],
                "notes": "Заметка докладчика про источник цифры",
            },
        ]
    )
    document = parse_document(data)
    assert document.kind == "pptx"
    assert document.anchor_unit == "слайд"
    assert [a.number for a in document.anchors] == [1, 2]
    assert "Заголовок первого слайда" in document.anchors[0].text
    second = document.anchors[1].text
    assert "Второй слайд" in second
    assert "1 200 млн" in second  # текст таблицы
    assert "Заметка докладчика" in second  # заметки докладчика


# --- parse_document: DOCX --------------------------------------------------


def test_parse_docx_numbers_non_empty_blocks_and_reads_tables():
    data = make_docx(
        ["Введение в тему", "", "   ", "Второй содержательный абзац"],
        table=[["Скважина", "Дебит"], ["№ 12", "48 т/сут"]],
    )
    document = parse_document(data)
    assert document.kind == "docx"
    assert document.anchor_unit == "блок"
    # 2 непустых абзаца + таблица одним блоком; пустые абзацы якорями не стали
    assert [a.number for a in document.anchors] == [1, 2, 3]
    assert document.anchors[0].text == "Введение в тему"
    assert document.anchors[1].text == "Второй содержательный абзац"
    assert "48 т/сут" in document.anchors[2].text


# --- parse_document: общее -------------------------------------------------


def test_parse_document_autodetects_kind():
    data = make_docx(["Автоопределение формата работает"])
    assert parse_document(data, kind=None).kind == "docx"


def test_parse_document_rejects_unknown_kind():
    with pytest.raises(UnsupportedFormatError):
        parse_document(make_pdf(["Hello there, a page of text"]), kind="xlsx")


def test_parse_document_rejects_garbage_bytes():
    with pytest.raises(UnsupportedFormatError):
        parse_document(b"\x00\x01\x02\x03\x04")


def test_parse_document_rejects_empty_bytes():
    with pytest.raises(UnsupportedFormatError):
        parse_document(b"")


# --- looks_scanned: границы порога -----------------------------------------


def _anchors(lengths: list[int]) -> list[Anchor]:
    return [Anchor(number=i, text="x" * n) for i, n in enumerate(lengths, start=1)]


def test_looks_scanned_on_empty_anchor_list():
    assert looks_scanned([], "pdf") is True


def test_looks_scanned_when_every_anchor_is_blank():
    assert looks_scanned(_anchors([0] * 5), "pdf") is True


def test_looks_scanned_on_page_number_artifacts():
    # Классический скан: pypdf достаёт с каждой страницы только колонцифру.
    assert looks_scanned(_anchors([3] * 12), "pdf") is True


def test_sparse_document_is_not_a_scan():
    # Три слайда по одной строке — редкий текст, но документ настоящий.
    assert looks_scanned(_anchors([30, 28, 35]), "pptx") is False


@pytest.mark.parametrize(
    "length, expected",
    [
        (SCAN_MEANINGFUL_ANCHOR_CHARS, False),  # ровно порог — уже содержательно
        (SCAN_MEANINGFUL_ANCHOR_CHARS - 1, True),  # на знак меньше — артефакт
    ],
)
def test_meaningful_anchor_threshold_both_sides(length, expected):
    assert looks_scanned(_anchors([length]), "pdf") is expected


@pytest.mark.parametrize(
    "filler, expected",
    [
        (5, False),  # плотность ровно 5.0 знаков на страницу — проходит
        (4, True),  # ниже порога — скан с редкими вкраплениями текста
    ],
)
def test_density_threshold_both_sides(filler, expected):
    # Один содержательный якорь среди сотни почти пустых первый критерий
    # проходит, и его ловит только плотность.
    anchors = _anchors([100] + [filler] * 99)
    total = sum(len(a.text) for a in anchors)
    assert (total / len(anchors) < SCAN_MIN_DENSITY) is expected
    assert looks_scanned(anchors, "pdf") is expected


def test_half_scanned_document_is_not_rejected():
    """Главный кейс калибровки: «Росатом» — 40 страниц с текстом из 107.

    Доля пустых якорей 63%, но текста 104 тысячи знаков. Критерий по доле
    пустых страниц выбросил бы этот документ целиком — он обязан разбираться.
    """
    anchors = _anchors([1000] * 40 + [0] * 67)
    assert looks_scanned(anchors, "pdf") is False


def test_density_rule_does_not_apply_to_docx():
    # У docx пустые блоки в якоря не попадают, мерить плотность там нечего.
    assert looks_scanned(_anchors([25] + [1] * 99), "docx") is False


# --- калибровка на реальных файлах (пропускается, если файлов нет) ---------

ROSATOM_PDF = "/Users/apple/Downloads/Росатом_позиция_+_заключение_экспертизы.pdf"
BPMN_PDF = "/Users/apple/Downloads/BPMN_2_0.pdf"
RADAR_PPTX = "/Users/apple/Downloads/Нефтесервисный радар.pptx"
NIR_DOCX = "/Users/apple/Downloads/Отчёт_НИР_Липскеров_МАС_нефтегаз_1.docx"


def _read(path: str) -> bytes:
    with open(path, "rb") as handle:
        return handle.read()


@pytest.mark.skipif(not os.path.exists(ROSATOM_PDF), reason="нет файла калибровки")
def test_calibration_half_scanned_pdf():
    data = _read(ROSATOM_PDF)
    assert detect_kind(data) == "pdf"
    document = parse_document(data)  # не должен упасть как «скан»
    assert len(document.anchors) == 107
    non_empty = [a for a in document.anchors if a.text]
    assert len(non_empty) == 40
    assert document.total_chars > 100_000
    assert looks_scanned(document.anchors, "pdf") is False


@pytest.mark.skipif(not os.path.exists(BPMN_PDF), reason="нет файла калибровки")
def test_calibration_full_text_pdf():
    document = parse_document(_read(BPMN_PDF))
    assert len(document.anchors) == 299
    assert all(a.text for a in document.anchors)
    assert document.total_chars > 600_000


@pytest.mark.skipif(not os.path.exists(RADAR_PPTX), reason="нет файла калибровки")
def test_calibration_real_pptx():
    data = _read(RADAR_PPTX)
    assert detect_kind(data) == "pptx"
    document = parse_document(data)
    assert len(document.anchors) == 13
    assert document.total_chars > 9_000
    assert [a.number for a in document.anchors] == list(range(1, 14))


@pytest.mark.skipif(not os.path.exists(NIR_DOCX), reason="нет файла калибровки")
def test_calibration_real_docx():
    data = _read(NIR_DOCX)
    assert detect_kind(data) == "docx"
    document = parse_document(data)
    # 90 непустых абзацев + 4 таблицы отдельными блоками
    assert len(document.anchors) == 94
    assert sum(1 for a in document.anchors if " | " in a.text) == 4
    assert document.total_chars > 20_000
